import os
import re
from dataclasses import dataclass
from typing import List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from sqlalchemy import Column, String, Integer, Date
from sqlalchemy.ext.declarative import declarative_base

Base = declarative_base()


class Deck(Base):
    __tablename__ = 'decks'

    id = Column(Integer, primary_key=True)
    title = Column(String)
    client = Column(String)
    date = Column(Date)
    category = Column(String)

    def __init__(self, title, client, date, category):
        self.title = title
        self.client = client
        self.date = date
        self.category = category


class Slide(Base):
    __tablename__ = 'slides'

    slide_id = Column(Integer, primary_key=True)
    text = Column(List)
    imgs = Column(List)
    deck_id = Column(Integer, ForeignKey('decks.id'))

    def __init__(self, slide_id, text, imgs, deck_id):
        self.slide_id = slide_id
        self.text = text
        self.imgs = imgs
        self.deck_id = deck_id


class Extractor:

    def __init__(self, file_path):
        self.file_path = file_path
        self.prs = Presentation(file_path)
        self.deck_name = self.set_deck_name()

    def set_deck_name(self):
        return re.findall(r'input_files/(.*?).pptx', self.file_path)[0]

    def write_image(self, shape, slide_id, shape_no):
        image = shape.image
        image_bytes = image.blob
        image_filename = f'{slide_id}{shape_no}.{image.ext}'
        image_file_path = f'images/{self.deck_name}/{image_filename}'
        with open(image_file_path, 'wb') as f:
            f.write(image_bytes)
        

    def create_img_dir(self):
        os.mkdir(f'images/{self.deck_name}')

    def extract_text_and_img(self):
        self.create_img_dir()
        deck_text_list = []
        for slide in self.prs.slides:
            slide_id = slide.slide_id
            slide_text_list = []
            for shape in slide.shapes:
                shape_no = 0
                if hasattr(shape, 'text'):
                    if shape.text != '':
                        slide_text_list.append(shape.text)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_path = self.write_image(shape, slide_id, shape_no)
                shape_no += 1
            deck_text_list.append(slide_text_list)

        return deck_text_list


output = Extractor('input_files/old_deck.pptx').extract_text_and_img()
print(output)
